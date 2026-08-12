"""prepare_receipts.scan_folder 단위 테스트. fitz/python-pptx로 최소 픽스처를 즉석 생성한다."""
import sys
import os
import tempfile
import shutil

SCRIPTS_DIR = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "scripts")
sys.path.insert(0, SCRIPTS_DIR)
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from test_support import run_test_classes
from PIL import Image
import fitz
from pptx import Presentation
from pptx.util import Inches

from prepare_receipts import scan_folder

LONG_PDF_TEXT = (
    "카드매출전표\n결제일시: 2026-09-01 12:30\n가맹점명: 테스트식당\n"
    "합계금액: 9,000원\n승인번호: 123456789\n부가세: 818원\n"
)


def _make_image_file(folder, name="receipt.png"):
    path = os.path.join(folder, name)
    Image.new("RGB", (200, 200), color="white").save(path)
    return path


def _make_text_pdf(folder, name="receipt.pdf"):
    path = os.path.join(folder, name)
    doc = fitz.open()
    page = doc.new_page()
    # fontname="korea": PyMuPDF's default "helv" (Helvetica/WinAnsi) cannot
    # encode Hangul and silently corrupts non-Latin text on extraction; the
    # bundled "korea" CJK font round-trips Korean text correctly.
    page.insert_text((72, 72), LONG_PDF_TEXT, fontname="korea")
    doc.save(path)
    doc.close()
    return path


def _make_pptx_file(folder, name="receipt.pptx"):
    img_path = _make_image_file(folder, "slide_img.png")
    path = os.path.join(folder, name)
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(img_path, Inches(1), Inches(1), height=Inches(3), width=Inches(3))
    prs.save(path)
    os.remove(img_path)
    return path


class TestScanFolder:
    def _run_in_temp_dirs(self, fn):
        src_dir = tempfile.mkdtemp(prefix="sikdae_src_")
        temp_dir = tempfile.mkdtemp(prefix="sikdae_tmp_")
        try:
            fn(src_dir, temp_dir)
        finally:
            shutil.rmtree(src_dir, ignore_errors=True)
            shutil.rmtree(temp_dir, ignore_errors=True)

    def test_finds_image_file(self):
        def check(src_dir, temp_dir):
            img_path = _make_image_file(src_dir)
            items = scan_folder(src_dir, temp_dir)
            assert len(items) == 1
            assert items[0]["source_file"] == img_path
            assert items[0]["image_path"] == img_path
        self._run_in_temp_dirs(check)

    def test_finds_pdf_with_embedded_text(self):
        def check(src_dir, temp_dir):
            _make_text_pdf(src_dir)
            items = scan_folder(src_dir, temp_dir)
            assert len(items) == 1
            assert "테스트식당" in items[0]["embedded_text"]
            assert os.path.exists(items[0]["image_path"])
        self._run_in_temp_dirs(check)

    def test_finds_pptx_image(self):
        def check(src_dir, temp_dir):
            pptx_path = _make_pptx_file(src_dir)
            items = scan_folder(src_dir, temp_dir)
            assert len(items) == 1
            assert items[0]["source_file"] == pptx_path
            assert os.path.exists(items[0]["image_path"])
        self._run_in_temp_dirs(check)

    def test_ignores_unsupported_files(self):
        def check(src_dir, temp_dir):
            with open(os.path.join(src_dir, "notes.txt"), "w") as f:
                f.write("무관한 파일")
            items = scan_folder(src_dir, temp_dir)
            assert len(items) == 0
        self._run_in_temp_dirs(check)

    def test_corrupt_pdf_recorded_as_error_without_aborting_scan(self):
        def check(src_dir, temp_dir):
            _make_image_file(src_dir, "good.png")
            bad_pdf_path = os.path.join(src_dir, "corrupt.pdf")
            with open(bad_pdf_path, "wb") as f:
                f.write(b"not a real pdf file")
            items = scan_folder(src_dir, temp_dir)
            assert len(items) == 2
            good_item = next(i for i in items if i["source_file"].endswith("good.png"))
            bad_item = next(i for i in items if i["source_file"] == bad_pdf_path)
            assert good_item["image_path"] is not None
            assert "error" in bad_item
            assert bad_item["image_path"] is None
        self._run_in_temp_dirs(check)

    def test_corrupt_pptx_recorded_as_error_without_aborting_scan(self):
        def check(src_dir, temp_dir):
            _make_image_file(src_dir, "good.png")
            bad_pptx_path = os.path.join(src_dir, "corrupt.pptx")
            with open(bad_pptx_path, "wb") as f:
                f.write(b"not a real pptx file")
            items = scan_folder(src_dir, temp_dir)
            assert len(items) == 2
            good_item = next(i for i in items if i["source_file"].endswith("good.png"))
            bad_item = next(i for i in items if i["source_file"] == bad_pptx_path)
            assert good_item["image_path"] is not None
            assert "error" in bad_item
            assert bad_item["image_path"] is None
        self._run_in_temp_dirs(check)


if __name__ == "__main__":
    success = run_test_classes(TestScanFolder)
    sys.exit(0 if success else 1)
