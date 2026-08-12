"""PPTX 파일에서 영수증 이미지를 추출하는 유틸리티."""

import io
from PIL import Image
from pptx import Presentation


def is_pptx(file_path: str) -> bool:
    return file_path.lower().endswith(".pptx")


def pptx_to_images(pptx_path: str) -> list[Image.Image]:
    """PPTX 파일에서 모든 이미지를 추출한다."""
    prs = Presentation(pptx_path)
    images = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == 13:  # Picture
                try:
                    blob = shape.image.blob
                    img = Image.open(io.BytesIO(blob))
                    # 너무 작은 이미지는 건너뛰기 (아이콘 등)
                    if img.width >= 100 and img.height >= 100:
                        if img.mode != "RGB":
                            img = img.convert("RGB")
                        images.append(img)
                except Exception:
                    continue

    return images
